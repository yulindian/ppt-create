#!/usr/bin/env python
"""Create a first-N-page PPTX subset for trial conversions."""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("input_pptx")
    parser.add_argument("output_pptx")
    parser.add_argument("--pages", type=int, default=3)
    args = parser.parse_args()

    if args.pages < 1:
        raise SystemExit("--pages must be positive")

    src = Path(args.input_pptx).resolve()
    out = Path(args.output_pptx).resolve()
    if not src.is_file():
        raise SystemExit(f"missing PPTX: {src}")

    prs = Presentation(str(src))
    keep = min(args.pages, len(prs.slides))
    sld_id_list = prs.slides._sldIdLst

    for idx in range(len(prs.slides) - 1, keep - 1, -1):
        r_id = sld_id_list[idx].rId
        prs.part.drop_rel(r_id)
        sld_id_list.remove(sld_id_list[idx])

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(out)
    print(f"source_slides={len(Presentation(str(src)).slides)}")
    print(f"subset_slides={len(Presentation(str(out)).slides)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

