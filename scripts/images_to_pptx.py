#!/usr/bin/env python3
import argparse
import json
import zipfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


SUPPORTED_SUFFIXES = {".png", ".jpg", ".jpeg"}


def collect_slide_images(slides_dir):
    files = [
        file for file in slides_dir.iterdir()
        if file.is_file()
        and file.name.lower().startswith("slide-")
        and file.suffix.lower() in SUPPORTED_SUFFIXES
    ]
    return sorted(files, key=lambda file: file.name.lower())


def add_full_slide_image(prs, image_path):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)


def main():
    parser = argparse.ArgumentParser(description="Create a PPTX from ordered slide images.")
    parser.add_argument("--slides-dir", required=True, help="Directory containing slide-01.png/jpg, slide-02.png/jpg, ...")
    parser.add_argument("--out", required=True, help="Output PPTX path")
    parser.add_argument("--expected-count", type=int, default=None, help="Fail if the image count differs")
    args = parser.parse_args()

    slides_dir = Path(args.slides_dir)
    out_pptx = Path(args.out)
    files = collect_slide_images(slides_dir)
    if not files:
        raise SystemExit(f"No slide images found in {slides_dir}")
    if args.expected_count is not None and len(files) != args.expected_count:
        raise SystemExit(f"Expected {args.expected_count} images, got {len(files)}")

    with Image.open(files[0]) as first:
        width_px, height_px = first.size
    ratio = width_px / height_px

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(13.333333 / ratio)

    # Remove the default empty first slide only if the template created one.
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    for file in files:
        add_full_slide_image(prs, file)

    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_pptx)

    with zipfile.ZipFile(out_pptx) as zf:
        slide_count = sum(1 for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml"))

    print(json.dumps({
        "output": str(out_pptx),
        "slides": slide_count,
        "images": len(files),
        "source_dir": str(slides_dir),
        "size_mb": round(out_pptx.stat().st_size / 1024 / 1024, 2),
    }, ensure_ascii=False))


if __name__ == "__main__":
    main()
